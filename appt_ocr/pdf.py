"""PDF to PPTX preprocessing.

Converts PDF pages to high-resolution images via PyMuPDF and wraps them
into a temporary PPTX file for the OCR pipeline.
"""

import io
import logging
import tempfile

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu

logger = logging.getLogger(__name__)


def convert_pdf_to_pptx(pdf_path: str, dpi: int = 300) -> str:
    """Convert a PDF file to PPTX.

    Renders each page as a high-resolution image and inserts it as a
    full-page slide in a new PPTX file.

    Args:
        pdf_path: Path to the PDF file.
        dpi: Rendering resolution (default 300, balancing quality and speed).

    Returns:
        File path of the generated temporary PPTX.
    """
    doc = fitz.open(pdf_path)
    prs = Presentation()

    # Set slide size based on the first page (PDF defaults to 72 DPI)
    first_page = doc[0]
    page_rect = first_page.rect
    slide_width_emu = int(page_rect.width * 914400 / 72)
    slide_height_emu = int(page_rect.height * 914400 / 72)
    prs.slide_width = slide_width_emu  # type: ignore[assignment]
    prs.slide_height = slide_height_emu  # type: ignore[assignment]

    # Use a blank layout
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]

        # Render page to PNG
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")

        # Dynamically adjust slide size per page
        page_rect = page.rect
        sw = int(page_rect.width * 914400 / 72)
        sh = int(page_rect.height * 914400 / 72)

        slide = prs.slides.add_slide(blank_layout)

        # Insert full-page picture
        slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            Emu(0),
            Emu(0),
            Emu(sw),
            Emu(sh),
        )

    total_pages = len(doc)
    doc.close()

    # Save as a temporary PPTX
    tmp_pptx = tempfile.NamedTemporaryFile(
        suffix=".pptx", delete=False, prefix="pdf2pptx_"
    )
    prs.save(tmp_pptx.name)
    tmp_pptx.close()

    logger.info("PDF conversion complete: %d pages -> %s", total_pages, tmp_pptx.name)
    return tmp_pptx.name
