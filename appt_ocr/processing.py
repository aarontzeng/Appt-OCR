"""Core slide and PPTX processing orchestration.

Coordinates the full pipeline: image extraction → OCR → text box merging →
regex filtering → inpainting → text box reconstruction.
"""

import io
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.shapes.picture import Picture
from pptx.util import Emu, Pt
from tqdm import tqdm

from appt_ocr.coordinates import compute_scale_factors, estimate_font_size, px_to_emu
from appt_ocr.image import analyze_text_features, extract_image_from_shape
from appt_ocr.inpainting import erase_text_using_masks
from appt_ocr.merging import merge_nearby_boxes
from appt_ocr.ocr import get_opencc_converter, run_ocr_on_image


def process_slide(
    slide,
    dpi: int = 96,
    lang: str = "ch",
    keep_images: bool = False,
    merge_threshold: float = 0.5,
    ignore_re: str = "",
    remove_re: str = "",
    inpaint_engine: str = "lama",
    watermark_only: bool = False,
    s2t: bool = False,
) -> int:
    """Process a single slide: Extract image → OCR → Add text boxes.

    Iterates through all Picture Shapes on the slide, performs OCR on each
    image, and restores text boxes based on the coordinates. Optionally
    deletes the original image after processing.

    Args:
        slide: python-pptx Slide object.
        dpi: Image DPI setting.
        lang: OCR language.
        keep_images: Whether to keep the original image.
        merge_threshold: Text box merging threshold.
        ignore_re: Regex for text to keep in background (no text box).
        remove_re: Regex for text to erase silently (no text box).
        inpaint_engine: ``"lama"`` or ``"opencv"``.
        watermark_only: If True, only erase text matching ``remove_re``.
        s2t: If True, convert simplified to traditional Chinese.

    Returns:
        Total number of text boxes identified on this slide.
    """
    # Collect all Picture Shapes first to avoid modifying during iteration
    picture_shapes: list[Picture] = [
        shape for shape in slide.shapes if isinstance(shape, Picture)
    ]

    total_boxes = 0

    for pic_shape in picture_shapes:
        # 1. Extract image
        image_blob, img_w_px, img_h_px = extract_image_from_shape(pic_shape)

        # 2. Calculate scaling factors
        shape_left_emu = pic_shape.left
        shape_top_emu = pic_shape.top
        shape_width_emu = pic_shape.width
        shape_height_emu = pic_shape.height

        scale_x, scale_y = compute_scale_factors(
            img_w_px,
            img_h_px,
            shape_width_emu,
            shape_height_emu,
            dpi,
        )

        # 3. Execute OCR
        ocr_results = run_ocr_on_image(image_blob, lang)
        if not ocr_results:
            continue

        # 4. Merge adjacent boxes
        merged_results = merge_nearby_boxes(ocr_results, merge_threshold)

        # 4.5. Filter text boxes using regular expressions
        kept_boxes = []  # Items requiring a newly created text box
        inpaint_boxes = []  # Items that need to be patched and erased

        for item in merged_results:
            text = item["text"].strip()

            # Detect rotated, vertical, or wrongly merged text boxes
            w, h = item["width_px"], item["height_px"]
            # (a) Obvious vertical text: height > 1.5x width -> skip
            if h > w * 1.5 and h > 50:
                continue
            # (b) Abnormal text density check
            est_text_width = len(text) * h * 0.55 * 0.72
            if w > 0 and est_text_width / w > 2.5:
                continue

            if remove_re and re.search(remove_re, text):
                # Matches remove_re (e.g. watermark): erase, no text box
                pad_x = int(item["width_px"] * 0.2)
                pad_y = int(item["height_px"] * 0.15)
                item["left_px"] = max(0, item["left_px"] - pad_x)
                item["top_px"] = max(0, item["top_px"] - pad_y)
                item["width_px"] += pad_x * 2
                item["height_px"] += pad_y * 2
                inpaint_boxes.append(item)
            elif watermark_only:
                # Watermark-only mode: skip all non-matching text
                pass
            elif ignore_re and re.search(ignore_re, text):
                # Matches ignore_re (e.g. math formula): keep in background
                pass
            else:
                # Regular text: erase AND create text box
                kept_boxes.append(item)
                inpaint_boxes.append(item)

        # 5. Analyze text features (color, boldness, and precise masks)
        for item in kept_boxes + inpaint_boxes:
            c_rgb, is_bold, mask = analyze_text_features(image_blob, item)
            item["color"] = c_rgb
            item["is_bold"] = is_bold
            item["text_mask"] = mask

        # 6. Erase text using inpainting and replace the original image
        if not keep_images and inpaint_boxes:
            inpainted_bytes = erase_text_using_masks(
                image_blob,
                inpaint_boxes,
                engine=inpaint_engine,
            )

            # Replace the image layer
            slide.shapes.add_picture(
                io.BytesIO(inpainted_bytes),
                shape_left_emu,
                shape_top_emu,
                shape_width_emu,
                shape_height_emu,
            )
            # Delete the old image layer
            sp = pic_shape._element
            sp.getparent().remove(sp)

        # 7. Convert OCR results to text boxes
        for item in kept_boxes:
            left_emu = (
                int(px_to_emu(item["left_px"], dpi) * scale_x) + shape_left_emu
            )
            top_emu = (
                int(px_to_emu(item["top_px"], dpi) * scale_y) + shape_top_emu
            )
            width_emu = int(px_to_emu(item["width_px"], dpi) * scale_x)
            height_emu = int(px_to_emu(item["height_px"], dpi) * scale_y)

            # Add horizontal padding to prevent wrapping
            width_emu = int(width_emu * 1.15)
            width_emu = max(width_emu, Emu(1))
            height_emu = max(height_emu, Emu(1))

            # Add text box
            txbox = slide.shapes.add_textbox(
                Emu(left_emu),
                Emu(top_emu),
                Emu(width_emu),
                Emu(height_emu),
            )
            tf = txbox.text_frame
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0

            # Set text content (with optional S2T conversion)
            text_content = item["text"]
            if s2t:
                cc = get_opencc_converter()
                if cc:
                    text_content = cc.convert(text_content)
            p = tf.paragraphs[0]
            p.text = text_content
            p.alignment = PP_ALIGN.LEFT

            # Dynamic font size and style
            font_pt = estimate_font_size(item["height_px"], scale_y, dpi)
            r, g, b = item["color"]
            is_bold = item.get("is_bold", False)

            for run in p.runs:
                run.font.size = Pt(font_pt)
                run.font.color.rgb = RGBColor(r, g, b)
                run.font.bold = is_bold

            total_boxes += 1

    return total_boxes


def process_pptx(
    input_path: str,
    output_path: str,
    dpi: int = 96,
    lang: str = "ch",
    keep_images: bool = False,
    merge_threshold: float = 0.5,
    ignore_re: str = "",
    remove_re: str = "",
    inpaint_engine: str = "lama",
    watermark_only: bool = False,
    s2t: bool = False,
) -> dict:
    """Process an entire PPTX file: Extract images → OCR → Restore text boxes.

    Args:
        input_path: Input PPTX file path.
        output_path: Output PPTX file path.
        dpi: Image DPI.
        lang: OCR language (``'ch'`` bilingual / ``'en'`` English only).
        keep_images: Whether to keep original image Shapes.
        merge_threshold: Text box merging threshold.
        ignore_re: Regex for text to keep in background.
        remove_re: Regex for text to erase silently.
        inpaint_engine: ``"lama"`` or ``"opencv"``.
        watermark_only: If True, only erase text matching ``remove_re``.
        s2t: If True, convert simplified to traditional Chinese.

    Returns:
        Processing statistics dict with keys: ``input``, ``output``,
        ``total_slides``, ``processed_slides``, ``total_textboxes``.
    """
    prs = Presentation(input_path)

    total_slides = len(prs.slides)
    processed_slides = 0
    total_textboxes = 0

    # Process page by page
    for _i, slide in enumerate(
        tqdm(
            prs.slides,
            desc=f"  Processing {Path(input_path).name}",
            unit="page",
        ),
    ):
        num_boxes = process_slide(
            slide,
            dpi=dpi,
            lang=lang,
            keep_images=keep_images,
            merge_threshold=merge_threshold,
            ignore_re=ignore_re,
            remove_re=remove_re,
            inpaint_engine=inpaint_engine,
            watermark_only=watermark_only,
            s2t=s2t,
        )
        if num_boxes > 0:
            processed_slides += 1
        total_textboxes += num_boxes

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    prs.save(output_path)

    return {
        "input": input_path,
        "output": output_path,
        "total_slides": total_slides,
        "processed_slides": processed_slides,
        "total_textboxes": total_textboxes,
    }
