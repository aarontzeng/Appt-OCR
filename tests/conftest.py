"""pytest configuration and fixtures."""

import io
import os
import tempfile
from pathlib import Path

import pytest
from PIL import Image


@pytest.fixture
def sample_image_bytes():
    """Create a sample image (RGB, 100x100) for testing."""
    img = Image.new("RGB", (100, 100), color=(255, 0, 0))
    img_bytes = io.BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)
    return img_bytes.getvalue()


@pytest.fixture
def sample_ocr_result():
    """Create a sample OCR result dictionary."""
    return {
        "left_px": 10.0,
        "top_px": 20.0,
        "width_px": 80.0,
        "height_px": 25.0,
        "text": "Hello World",
        "confidence": 0.95,
    }


@pytest.fixture
def sample_ocr_results():
    """Create sample OCR results with multiple boxes."""
    return [
        {
            "left_px": 10.0,
            "top_px": 20.0,
            "width_px": 80.0,
            "height_px": 25.0,
            "text": "Hello",
            "confidence": 0.95,
        },
        {
            "left_px": 100.0,
            "top_px": 20.0,
            "width_px": 80.0,
            "height_px": 25.0,
            "text": "World",
            "confidence": 0.92,
        },
    ]


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield tmpdir


@pytest.fixture
def sample_pptx_path(temp_dir):
    """Create a simple sample PPTX file for testing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Inches(1)
        pic_path = os.path.join(temp_dir, "test_image.png")

        # Create a test image
        img = Image.new("RGB", (200, 150), color=(100, 150, 200))
        img.save(pic_path)

        # Add image to slide
        slide.shapes.add_picture(pic_path, left, top, width=Inches(4))

        # Save PPTX
        pptx_path = os.path.join(temp_dir, "test.pptx")
        prs.save(pptx_path)
        return pptx_path
    except Exception as e:
        pytest.skip(f"Could not create sample PPTX: {e}")


@pytest.fixture
def sample_pdf_path(temp_dir):
    """Create a simple sample PDF file for testing."""
    try:
        import fitz

        doc = fitz.open()
        page = doc.new_page(width=200, height=150)
        page.insert_text((10, 10), "Test PDF Content", fontsize=12)
        pdf_path = os.path.join(temp_dir, "test.pdf")
        doc.save(pdf_path)
        return pdf_path
    except ImportError:
        pytest.skip("PyMuPDF not available")
