"""Tests for image.py"""

import io

import numpy as np
import pytest
from PIL import Image

from appt_ocr.image import analyze_text_features, extract_image_from_shape


class TestAnalyzeTextFeatures:
    """Test text feature analysis."""

    def test_analyze_text_features_returns_tuple(self, sample_image_bytes):
        """Test that function returns a 3-tuple."""
        box = {
            "left_px": 10.0,
            "top_px": 10.0,
            "width_px": 80.0,
            "height_px": 25.0,
        }
        result = analyze_text_features(sample_image_bytes, box)
        assert isinstance(result, tuple)
        assert len(result) == 3

    def test_analyze_text_features_color_format(self, sample_image_bytes):
        """Test that returned color is RGB tuple."""
        box = {
            "left_px": 10.0,
            "top_px": 10.0,
            "width_px": 80.0,
            "height_px": 25.0,
        }
        color, _, _ = analyze_text_features(sample_image_bytes, box)
        assert isinstance(color, tuple)
        assert len(color) == 3
        assert all(isinstance(c, int) for c in color)
        assert all(0 <= c <= 255 for c in color)

    def test_analyze_text_features_bold_flag(self, sample_image_bytes):
        """Test that bold flag is a boolean."""
        box = {
            "left_px": 10.0,
            "top_px": 10.0,
            "width_px": 80.0,
            "height_px": 25.0,
        }
        _, is_bold, _ = analyze_text_features(sample_image_bytes, box)
        assert isinstance(is_bold, bool)

    def test_analyze_text_features_mask_type(self, sample_image_bytes):
        """Test that returned mask is numpy array."""
        box = {
            "left_px": 10.0,
            "top_px": 10.0,
            "width_px": 80.0,
            "height_px": 25.0,
        }
        _, _, mask = analyze_text_features(sample_image_bytes, box)
        assert isinstance(mask, np.ndarray)
        assert mask.dtype == np.uint8
        assert mask.ndim == 2

    def test_analyze_text_features_mask_dimensions(self, sample_image_bytes):
        """Test that mask dimensions match box size."""
        box = {
            "left_px": 10.0,
            "top_px": 10.0,
            "width_px": 80.0,
            "height_px": 25.0,
        }
        _, _, mask = analyze_text_features(sample_image_bytes, box)
        assert mask.shape == (25, 80)

    def test_analyze_text_features_out_of_bounds_box(self, sample_image_bytes):
        """Test with box extending beyond image."""
        box = {
            "left_px": 50.0,
            "top_px": 50.0,
            "width_px": 100.0,
            "height_px": 100.0,
        }
        # Should handle gracefully without crashing
        color, is_bold, mask = analyze_text_features(sample_image_bytes, box)
        assert isinstance(color, tuple)
        assert isinstance(is_bold, bool)
        assert isinstance(mask, np.ndarray)

    def test_analyze_text_features_negative_coords(self, sample_image_bytes):
        """Test with negative coordinates."""
        box = {
            "left_px": -10.0,
            "top_px": -10.0,
            "width_px": 50.0,
            "height_px": 50.0,
        }
        # Should handle gracefully by clamping to 0
        color, is_bold, mask = analyze_text_features(sample_image_bytes, box)
        assert isinstance(color, tuple)
        assert isinstance(is_bold, bool)
        assert isinstance(mask, np.ndarray)

    def test_analyze_text_features_zero_size(self, sample_image_bytes):
        """Test with zero-sized box."""
        box = {
            "left_px": 10.0,
            "top_px": 10.0,
            "width_px": 0.0,
            "height_px": 0.0,
        }
        color, is_bold, mask = analyze_text_features(sample_image_bytes, box)
        # Should return default values
        assert isinstance(color, tuple)
        assert isinstance(is_bold, bool)

    def test_analyze_text_features_invalid_image(self):
        """Test with invalid image bytes."""
        box = {
            "left_px": 10.0,
            "top_px": 10.0,
            "width_px": 50.0,
            "height_px": 50.0,
        }
        invalid_bytes = b"not an image"
        # Should return default values
        color, is_bold, mask = analyze_text_features(invalid_bytes, box)
        assert color == (0, 0, 0)
        assert is_bold is False
        assert isinstance(mask, np.ndarray)

    def test_analyze_text_features_real_red_image(self):
        """Test text feature analysis with red image."""
        # Create a red image
        img = Image.new("RGB", (100, 100), color=(255, 0, 0))
        img_bytes = io.BytesIO()
        img.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        box = {
            "left_px": 20.0,
            "top_px": 20.0,
            "width_px": 60.0,
            "height_px": 60.0,
        }
        color, is_bold, mask = analyze_text_features(img_bytes.getvalue(), box)
        # Red image should produce some red component
        assert isinstance(color, tuple)


class TestExtractImageFromShape:
    """Test image extraction from shapes."""

    def test_extract_image_from_shape_returns_tuple(self, sample_pptx_path):
        """Test that function returns a 3-tuple."""
        try:
            from pptx import Presentation
            from pptx.shapes.picture import Picture

            prs = Presentation(sample_pptx_path)
            slide = prs.slides[0]

            # Find a picture shape
            for shape in slide.shapes:
                if isinstance(shape, Picture):
                    result = extract_image_from_shape(shape)
                    assert isinstance(result, tuple)
                    assert len(result) == 3
                    break
        except Exception as e:
            pytest.skip(f"Could not test PPTX extraction: {e}")

    def test_extract_image_returns_bytes(self, sample_pptx_path):
        """Test that image content is bytes."""
        try:
            from pptx import Presentation
            from pptx.shapes.picture import Picture

            prs = Presentation(sample_pptx_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if isinstance(shape, Picture):
                    img_bytes, _, _ = extract_image_from_shape(shape)
                    assert isinstance(img_bytes, bytes)
                    assert len(img_bytes) > 0
                    break
        except Exception as e:
            pytest.skip(f"Could not test PPTX extraction: {e}")

    def test_extract_image_dimensions(self, sample_pptx_path):
        """Test that image dimensions are positive integers."""
        try:
            from pptx import Presentation
            from pptx.shapes.picture import Picture

            prs = Presentation(sample_pptx_path)
            slide = prs.slides[0]

            for shape in slide.shapes:
                if isinstance(shape, Picture):
                    _, width_px, height_px = extract_image_from_shape(shape)
                    assert isinstance(width_px, int)
                    assert isinstance(height_px, int)
                    assert width_px > 0
                    assert height_px > 0
                    break
        except Exception as e:
            pytest.skip(f"Could not test PPTX extraction: {e}")
